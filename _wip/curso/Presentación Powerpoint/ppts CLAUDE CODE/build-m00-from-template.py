#!/usr/bin/env python3
"""
build-m00-from-template.py — Build M00 v2 from M04-PATRON.pptx template.

Reads: M04-PATRON.pptx (template), M00-El-Momento-que-Nos-Convoca.pptx (notes)
Output: M00-El-Momento-que-Nos-Convoca-v2.pptx

Usage: python build-m00-from-template.py
"""

import copy
import sys
from io import BytesIO
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 1: CONSTANTS
# ═══════════════════════════════════════════════════════════════════════════

SCRIPT_DIR = Path(__file__).parent
TEMPLATE_FILE = SCRIPT_DIR / "M04-PATRON.pptx"
M00_V1_FILE = SCRIPT_DIR / "M00-El-Momento-que-Nos-Convoca.pptx"
OUTPUT_FILE = SCRIPT_DIR / "M00-El-Momento-que-Nos-Convoca-v2.pptx"

SLIDE_W = 9144000
SLIDE_H = 5143500

# Colors (from M07 design system)
CORAL    = RGBColor(0xDC, 0x80, 0x60)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE   = RGBColor(0xE0, 0xEB, 0xEB)
META     = RGBColor(0xA8, 0xB0, 0xA8)
PETROL   = RGBColor(0x2A, 0x50, 0x58)
BODY     = RGBColor(0x5A, 0x6A, 0x6A)
CAPTION  = RGBColor(0x88, 0x90, 0x88)
CARD_BG  = RGBColor(0xF5, 0xF5, 0xF0)

FONT = "Montserrat"

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Template slide indices (0-based)
T_PORTADA    = 0
T_SEPARATOR  = 1
T_CONTRAST   = 2
T_PROCESS    = 3
T_GRID       = 7
T_CHART      = 8
T_VISUAL     = 9
T_CAUSE      = 10
T_CARDS      = 14
T_ROADMAP    = 17
T_TRANSITION = 18
T_PRACTICE   = 19


# ═══════════════════════════════════════════════════════════════════════════
# SECTION 2: LOW-LEVEL UTILITIES
# ═══════════════════════════════════════════════════════════════════════════

def clone_slide(prs, source_idx):
    """Clone slide at source_idx, append to end. Returns new slide."""
    source = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(source.slide_layout)

    # Clear default shapes
    spTree = new_slide.shapes._spTree
    for child in list(spTree):
        if etree.QName(child).localname not in ('nvGrpSpPr', 'grpSpPr'):
            spTree.remove(child)

    # Map image relationships from source to new slide
    rId_map = {}
    for key, rel in source.part.rels.items():
        if "image" in str(rel.reltype).lower():
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rId_map[key] = new_rId

    # Deep copy shapes, updating image rIds
    for child in source.shapes._spTree:
        if etree.QName(child).localname not in ('nvGrpSpPr', 'grpSpPr'):
            new_child = copy.deepcopy(child)
            for blip in new_child.iter(f'{{{NS_A}}}blip'):
                old_rId = blip.get(f'{{{NS_R}}}embed')
                if old_rId and old_rId in rId_map:
                    blip.set(f'{{{NS_R}}}embed', rId_map[old_rId])
            spTree.append(new_child)

    return new_slide


def remove_slides(prs, indices):
    """Remove slides at given 0-based indices."""
    sldIdLst = prs.slides._sldIdLst
    sldIds = list(sldIdLst)
    for idx in sorted(indices, reverse=True):
        sldId = sldIds[idx]
        rId = sldId.get(f'{{{NS_R}}}id')
        sldIdLst.remove(sldId)
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass


def clear_shapes(slide):
    """Remove all shapes from slide and clean up orphaned image rels."""
    spTree = slide.shapes._spTree
    orphan_rIds = set()
    for child in list(spTree):
        tag = etree.QName(child).localname
        if tag == "pic":
            blip = child.find(f".//{{{NS_A}}}blip")
            if blip is not None:
                rId = blip.get(f"{{{NS_R}}}embed")
                if rId:
                    orphan_rIds.add(rId)
    for child in list(spTree):
        if etree.QName(child).localname in ("sp", "pic", "grpSp", "cxnSp"):
            spTree.remove(child)
    for rId in orphan_rIds:
        try:
            slide.part.drop_rel(rId)
        except (KeyError, ValueError):
            pass


def extract_bg(slide):
    """Extract full-slide background image bytes."""
    for shape in slide.shapes:
        if shape.shape_type == 13 and shape.left == 0 and shape.top == 0:
            return shape.image.blob
    return None


def add_bg(slide, img_bytes, prs):
    """Add full-slide background image as first shape."""
    pic = slide.shapes.add_picture(BytesIO(img_bytes), 0, 0,
                                   prs.slide_width, prs.slide_height)
    spTree = slide.shapes._spTree
    el = pic._element
    spTree.remove(el)
    spTree.insert(2, el)


def tb(slide, left, top, w, h, text, pt, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    """Add textbox (shorthand). Returns shape."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(pt)
        run.font.bold = bold
        run.font.color.rgb = color
        if italic:
            run.font.italic = True
    return txBox


def hline(slide, left, top, width, color=PETROL):
    """Add thin horizontal line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(14288))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def card_bg(slide, left, top, w, h):
    """Add card rectangle with top accent bar."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(w), Emu(h))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.fill.background()
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(w), Emu(42863))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CORAL
    bar.line.fill.background()


def header(slide, label, title):
    """Standard content header (label + title) for light-bg slides."""
    tb(slide, 571500, 428625, 8001000, 142875,
       label.upper(), 10, bold=True, color=CORAL)
    tb(slide, 571500, 628650, 8001000, 285750,
       title, 22, bold=True, color=PETROL)


def quote_bar(slide, left, top, w, text, pt=9):
    """Quote box with left accent bar."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(28575), Emu(400000))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CORAL
    bar.line.fill.background()
    tb(slide, left + 100000, top + 50000, w - 100000, 350000,
       f"\u201C{text}\u201D", pt, italic=True, color=BODY)


def Q(t):
    """Wrap text in curly quotes."""
    return f"\u201C{t}\u201D"


# ═══════════════════════════════════════════════════════════════════════════
# SECTION 3: PATTERN BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

def build_portada(slide, bg, prs, etiqueta, titulo, cita):
    """Pattern #1: Portada de modulo (dark bg)."""
    clear_shapes(slide)
    add_bg(slide, bg, prs)
    tb(slide, 714375, 1006041, 6786563, 233958,
       etiqueta, 12, bold=True, color=CORAL)
    tb(slide, 714375, 1454311, 6786563, 1200150,
       titulo, 41, bold=True, color=WHITE)
    tb(slide, 1035844, 3497424, 6465094, 640035,
       Q(cita), 17, color=SUBTLE)


def build_separador(slide, bg, prs, num, titulo, sub, label="TEMA"):
    """Pattern #2: Separador de seccion (dark bg)."""
    clear_shapes(slide)
    add_bg(slide, bg, prs)
    tb(slide, 714375, 1617166, 1471054, 1028700,
       num, 93, bold=True, color=CORAL)
    tb(slide, 3042679, 1688604, 4920816, 194667,
       label.upper(), 10, bold=True, color=SUBTLE)
    tb(slide, 3042679, 2026146, 4920816, 1005818,
       titulo, 33, bold=True, color=WHITE)
    tb(slide, 3042679, 3246276, 4920816, 280029,
       sub, 15, color=SUBTLE)


def build_contraste(slide, etiq, titulo, lh, li, rh, ri, cita=None, sym=None):
    """Pattern #3: Contraste dual (light bg)."""
    clear_shapes(slide)
    header(slide, etiq, titulo)
    # Left column
    tb(slide, 571500, 1371600, 3714750, 357188, lh, 12, bold=True, color=PETROL)
    hline(slide, 501161, 1690871, 2926666)
    y = 2077380
    for item in li:
        tb(slide, 669727, y, 89297, 200025, "\u2022", 9, bold=True, color=CORAL)
        tb(slide, 1000125, y, 2700000, 188612, item, 9, color=BODY)
        y += 550068
    # Right column
    tb(slide, 4857750, 1371600, 3714750, 371475, rh, 12, bold=True, color=PETROL)
    hline(slide, 4857750, 1690871, 2926666)
    y = 2097416
    for item in ri:
        tb(slide, 4914900, y, 171450, 228600, "\u2713", 9, bold=True, color=CORAL)
        tb(slide, 5286375, y, 2700000, 205718, item, 9, color=BODY)
        y += 530000
    if sym:
        tb(slide, 4200000, 2400000, 400000, 400000, sym, 36, bold=True,
           color=META, align=PP_ALIGN.CENTER)
    if cita:
        quote_bar(slide, 571500, 4300000, 8001000, cita)


def build_tarjetas(slide, etiq, titulo, cards):
    """Pattern #15: Tarjetas de sintesis (light bg). 3 or 4 cards."""
    clear_shapes(slide)
    header(slide, etiq, titulo)
    n = len(cards)
    gap = 143259
    total_w = 8001000
    cw = (total_w - (n - 1) * gap) // n
    sx = 571500
    ct = 1200150
    ch = 3514725
    for i, c in enumerate(cards):
        x = sx + i * (cw + gap)
        card_bg(slide, x, ct, cw, ch)
        ix = x + 214313
        iw = cw - 428626
        num = c.get("num", f"{i+1:02d}")
        tb(slide, x + cw - 742613, ct + 214313, 528526, 457200,
           num, 36, bold=True, color=META, align=PP_ALIGN.RIGHT)
        if c.get("header"):
            tb(slide, ix, ct + 814388, iw, 457200,
               c["header"], 12, bold=True, color=PETROL)
        if c.get("subtitle"):
            tb(slide, ix, ct + 1200000, iw, 200025,
               Q(c["subtitle"]), 9, italic=True, color=BODY)
        desc_top = ct + 1500000 if c.get("subtitle") else ct + 1385888
        if c.get("text"):
            tb(slide, ix, desc_top, iw, 928688, c["text"], 9, color=BODY)


def build_proceso(slide, etiq, titulo, fases, placeholder=None):
    """Pattern #4: Proceso secuencial (light bg)."""
    clear_shapes(slide)
    header(slide, etiq, titulo)
    n = len(fases)
    gap = 57150
    total_w = 8001000
    cw = (total_w - (n - 1) * gap) // n
    sx = 571500
    ct = 1200150
    ch = 2743200
    for i, f in enumerate(fases):
        x = sx + i * (cw + gap)
        card_bg(slide, x, ct, cw, ch)
        ix = x + 150000
        iw = cw - 300000
        tb(slide, x + cw - 478463, ct + 71438, 400000, 685800,
           str(f["num"]), 54, bold=True, color=META, align=PP_ALIGN.RIGHT)
        tb(slide, ix, ct + 1000000, iw, 228600,
           f["header"], 12, bold=True, color=PETROL)
        tb(slide, ix, ct + 1400000, iw, 557213,
           f["text"], 9, color=BODY)
    if placeholder:
        tb(slide, 571500, 4100000, 8001000, 200000,
           f"[IMAGEN: {placeholder}]", 8, color=META, align=PP_ALIGN.CENTER)


def build_grid_cols(slide, etiq, titulo, items, cita=None):
    """Pattern #8 variant: Grid as horizontal columns (light bg)."""
    clear_shapes(slide)
    header(slide, etiq, titulo)
    n = len(items)
    gap = 100000
    total_w = 8001000
    cw = (total_w - (n - 1) * gap) // n
    sx = 571500
    for i, item in enumerate(items):
        x = sx + i * (cw + gap)
        tb(slide, x, 1200150, cw, 228600,
           item["header"], 11, bold=True, color=PETROL)
        tb(slide, x, 1500000, cw, 400000, item["text"], 9, color=BODY)
        if "stat" in item:
            tb(slide, x, 2000000, cw, 300000,
               item["stat"], 20, bold=True, color=CORAL)
    if cita:
        quote_bar(slide, 571500, 4300000, 8001000, cita)


def build_grid_rows(slide, etiq, titulo, items, cita=None):
    """Pattern #8 variant: Grid as vertical rows (light bg)."""
    clear_shapes(slide)
    header(slide, etiq, titulo)
    y = 1200150
    rh = 550000
    for item in items:
        tb(slide, 571500, y, 2500000, 200025,
           item["header"], 11, bold=True, color=PETROL)
        tb(slide, 3200000, y, 5372500, 200025,
           item["text"], 9, color=BODY)
        y += rh
    if cita:
        quote_bar(slide, 571500, min(y + 200000, 4400000), 8001000, cita)


def build_stat(slide, stat, desc, fuente=None, apoyo=None, placeholder=None):
    """Pattern #19: Stat destacado (light bg)."""
    clear_shapes(slide)
    tb(slide, 571500, 1200000, 4000000, 900000,
       stat, 72, bold=True, color=CORAL)
    tb(slide, 571500, 2100000, 4000000, 400000, desc, 14, color=PETROL)
    if placeholder:
        tb(slide, 5000000, 1200000, 3572500, 2000000,
           f"[IMAGEN:\n{placeholder}]", 10, color=META, align=PP_ALIGN.CENTER)
    if apoyo:
        tb(slide, 571500, 2700000, 7500000, 600000, apoyo, 9, color=BODY)
    if fuente:
        tb(slide, 571500, 4500000, 8001000, 200000,
           fuente, 8, italic=True, color=CAPTION)


def build_multi_stat(slide, etiq, titulo, stats, cita=None):
    """Pattern #19 variant: Multi-stat grid (light bg)."""
    clear_shapes(slide)
    header(slide, etiq, titulo)
    n = len(stats)
    gap = 150000
    total_w = 8001000
    cw = (total_w - (n - 1) * gap) // n
    sx = 571500
    for i, st in enumerate(stats):
        x = sx + i * (cw + gap)
        tb(slide, x, 1200000, cw, 600000,
           st["num"], 48, bold=True, color=CORAL)
        tb(slide, x, 1850000, cw, 400000, st["desc"], 9, color=BODY)
        if "source" in st:
            tb(slide, x, 2350000, cw, 200000,
               st["source"], 7, italic=True, color=CAPTION)
    if cita:
        quote_bar(slide, 571500, 3100000, 8001000, cita)


def build_cita(slide, texto, autor=None, obra=None, apoyo=None, fuente=None):
    """Pattern #18: Cita de autor (light bg)."""
    clear_shapes(slide)
    # Decorative quote mark
    tb(slide, 571500, 800000, 500000, 600000,
       "\u201C", 96, bold=True, color=CORAL)
    # Vertical accent line
    vline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(1200000), Emu(1200000),
        Emu(28575), Emu(1800000))
    vline.fill.solid()
    vline.fill.fore_color.rgb = CORAL
    vline.line.fill.background()
    # Quote text
    tb(slide, 1500000, 1300000, 6500000, 800000,
       Q(texto), 22, italic=True, color=PETROL)
    y = 2300000
    if autor:
        tb(slide, 1500000, y, 6500000, 250000,
           f"\u2014 {autor}", 12, bold=True, color=CORAL)
        y += 300000
    if obra:
        tb(slide, 1500000, y, 6500000, 200000,
           obra, 8, italic=True, color=CAPTION)
        y += 400000
    if apoyo:
        tb(slide, 1500000, y, 6500000, 800000, apoyo, 9, color=BODY)
        y += 900000
    if fuente:
        tb(slide, 1500000, min(y, 4200000), 6500000, 200000,
           fuente, 8, italic=True, color=CAPTION)


def build_reflexion(slide, bg, prs, pregunta, apoyo=None, opciones=None,
                    etiq=None):
    """Pattern #20: Reflexion / pregunta (dark bg)."""
    clear_shapes(slide)
    add_bg(slide, bg, prs)
    if etiq:
        tb(slide, 571500, 1000000, 8001000, 200000,
           etiq.upper(), 10, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    tb(slide, 571500, 1500000, 8001000, 800000,
       pregunta, 24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if apoyo:
        tb(slide, 571500, 2500000, 8001000, 300000,
           apoyo, 12, color=SUBTLE, align=PP_ALIGN.CENTER)
    if opciones:
        tb(slide, 571500, 3200000, 8001000, 300000,
           "  \u00B7  ".join(opciones), 11, color=CORAL,
           align=PP_ALIGN.CENTER)


def build_causa_efecto(slide, etiq, titulo, sub, filas, placeholder=None,
                       cita=None):
    """Pattern #11: Causa-efecto en filas (light bg)."""
    clear_shapes(slide)
    header(slide, etiq, titulo)
    if sub:
        tb(slide, 571500, 1000000, 8001000, 200000,
           sub, 9, italic=True, color=BODY)
    y = 1300000
    rh = 380000
    for i, fila in enumerate(filas):
        num = f"{i+1}."
        tb(slide, 571500, y, 300000, 200000, num, 9, bold=True, color=CORAL)
        tb(slide, 900000, y, 7672500, 250000, fila, 9, color=BODY)
        y += rh
    if placeholder:
        tb(slide, 571500, y + 100000, 8001000, 200000,
           f"[IMAGEN: {placeholder}]", 8, color=META, align=PP_ALIGN.CENTER)
        y += 400000
    if cita:
        quote_bar(slide, 571500, min(y + 100000, 4300000), 8001000, cita)


def build_visual_central(slide, etiq, titulo, placeholder, texto,
                         texto2=None):
    """Pattern #10: Visual central + bloques (light bg, simplified)."""
    clear_shapes(slide)
    header(slide, etiq, titulo)
    # Image placeholder (center)
    tb(slide, 2500000, 1200000, 4144000, 1500000,
       f"[IMAGEN: {placeholder}]", 10, color=META, align=PP_ALIGN.CENTER)
    # Text block below
    tb(slide, 571500, 2900000, 8001000, 800000, texto, 9, color=BODY)
    if texto2:
        tb(slide, 571500, 3800000, 8001000, 400000,
           texto2, 9, italic=True, color=BODY)


def build_grafica(slide, etiq, titulo, bloques, placeholder=None):
    """Pattern #9: Grafica + analisis (light bg)."""
    clear_shapes(slide)
    header(slide, etiq, titulo)
    if placeholder:
        tb(slide, 571500, 1000000, 8001000, 300000,
           f"[IMAGEN: {placeholder}]", 10, color=META, align=PP_ALIGN.CENTER)
    n = len(bloques)
    gap = 150000
    total_w = 8001000
    bw = (total_w - (n - 1) * gap) // n
    sx = 571500
    bt = 1600000
    for i, bloque in enumerate(bloques):
        x = sx + i * (bw + gap)
        card_bg(slide, x, bt, bw, 2600000)
        ix = x + 150000
        iw = bw - 300000
        tb(slide, ix, bt + 150000, iw, 228600,
           bloque["header"], 11, bold=True, color=PETROL)
        if bloque.get("subtitle"):
            tb(slide, ix, bt + 450000, iw, 200000,
               bloque["subtitle"], 9, italic=True, color=CORAL)
        tb(slide, ix, bt + 750000, iw, 1500000,
           bloque["text"], 9, color=BODY)


def build_roadmap(slide, etiq, titulo, items, cita=None):
    """Pattern #17: Roadmap / Timeline (light bg)."""
    clear_shapes(slide)
    header(slide, etiq, titulo)
    n = len(items)
    if n <= 4:
        # Horizontal layout
        gap = 100000
        total_w = 8001000
        bw = (total_w - (n - 1) * gap) // n
        sx = 571500
        for i, item in enumerate(items):
            x = sx + i * (bw + gap)
            card_bg(slide, x, 1200000, bw, 2500000)
            ix = x + 100000
            iw = bw - 200000
            tb(slide, ix, 1350000, iw, 228600,
               item["label"], 11, bold=True, color=PETROL)
            if item.get("subtitle"):
                tb(slide, ix, 1650000, iw, 200000,
                   item["subtitle"], 9, italic=True, color=CORAL)
            if item.get("text"):
                tb(slide, ix, 1950000, iw, 1200000,
                   item["text"], 9, color=BODY)
    else:
        # Vertical list for many items
        y = 1200000
        rh = 370000
        for item in items:
            tb(slide, 571500, y, 1500000, 200000,
               item["label"], 11, bold=True, color=PETROL)
            subtitle = item.get("subtitle", "")
            if subtitle:
                tb(slide, 2200000, y, 2000000, 200000,
                   subtitle, 9, italic=True, color=CORAL)
            text = item.get("text", "")
            if text:
                tb(slide, 4400000, y, 4172500, 200000,
                   text, 9, color=BODY)
            y += rh
    if cita:
        quote_bar(slide, 571500, 4400000, 8001000, cita)


def build_transicion(slide, bg, prs, modulo, nombre, cita,
                     placeholder=None):
    """Pattern #22: Transicion al siguiente modulo (dark bg)."""
    clear_shapes(slide)
    add_bg(slide, bg, prs)
    tb(slide, 571500, 910828, 8001000, 200025,
       "PROXIMA PARADA", 10, bold=True, color=SUBTLE,
       align=PP_ALIGN.CENTER)
    tb(slide, 571500, 1953816, 8001000, 685800,
       modulo, 54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, 571500, 2753916, 8001000, 342900,
       nombre, 24, color=SUBTLE, align=PP_ALIGN.CENTER)
    if placeholder:
        tb(slide, 3500000, 1300000, 2144000, 500000,
           f"[IMAGEN: {placeholder}]", 8, color=META,
           align=PP_ALIGN.CENTER)
    hline(slide, 2638611, 3439716, 3866750, SUBTLE)
    tb(slide, 2000000, 3671888, 5144000, 557213,
       Q(cita), 12, italic=True, color=SUBTLE, align=PP_ALIGN.CENTER)


def build_cierre(slide, bg, prs):
    """Pattern #21: Cierre / Gracias (dark bg)."""
    clear_shapes(slide)
    add_bg(slide, bg, prs)
    tb(slide, 2594883, 1456432, 3954205, 857250,
       "Gracias", 62, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    tb(slide, 2594883, 2770882, 3954205, 312539,
       "Liderar con Consciencia", 17, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, 2594883, 3512046, 3954205, 175022,
       "www.liderarconconsciencia.com", 9, color=META,
       align=PP_ALIGN.CENTER)


def build_sep_practica(slide, bg, prs, titulo, sub, detalle=None):
    """Pattern #2 variant: Separador de practica (dark bg)."""
    clear_shapes(slide)
    add_bg(slide, bg, prs)
    tb(slide, 3000375, 1653778, 4921151, 200025,
       "PRACTICA", 10, bold=True, color=CORAL)
    tb(slide, 3000375, 1968103, 4921151, 1071563,
       titulo, 33, bold=True, color=WHITE)
    tb(slide, 3000375, 3211116, 4921151, 278606,
       sub, 12, color=SUBTLE)
    if detalle:
        tb(slide, 3000375, 3600000, 4921151, 200000,
           detalle, 9, color=META)


# ═══════════════════════════════════════════════════════════════════════════
# SECTION 4: MAIN PIPELINE
# ═══════════════════════════════════════════════════════════════════════════

def build_m00():
    """Build all 30 slides of M00 v2."""
    print("Loading template and M00 v1...")
    prs = Presentation(str(TEMPLATE_FILE))
    m00v1 = Presentation(str(M00_V1_FILE))
    num_template = len(prs.slides)
    print(f"  Template: {num_template} slides")
    print(f"  M00 v1: {len(m00v1.slides)} slides")

    # Extract dark backgrounds from template
    bg_portada = extract_bg(prs.slides[T_PORTADA])
    bg_sep = extract_bg(prs.slides[T_SEPARATOR])
    bg_trans = extract_bg(prs.slides[T_TRANSITION])
    bg_pract = extract_bg(prs.slides[T_PRACTICE])

    print("Building 30 slides...")

    # For each M00 slide: clone a template slide, then rebuild content.
    # We clone from the source that has the right slide layout.
    # Dark-bg slides: clone from separator (has dark bg image)
    # Light-bg slides: clone from cards (clean content layout)
    slides = []

    # ── S1: Portada ───────────────────────────────────────────────────
    s = clone_slide(prs, T_PORTADA)
    build_portada(s, bg_portada, prs,
                  "Modulo 00",
                  "El Momento que\nNos Convoca",
                  "La consciencia se entrena, no se ensena")
    slides.append(("S1 Portada", s))

    # ── S2: Hoja de ruta (4 tarjetas) ────────────────────────────────
    s = clone_slide(prs, T_CARDS)
    build_tarjetas(s, "HOJA DE RUTA", "Estructura del modulo", [
        {"num": "01", "header": "Diagnostico",
         "text": "La tormenta perfecta: cuatro fuerzas que convergen."},
        {"num": "02", "header": "Transicion",
         "text": "La tentacion de la impotencia y las respuestas naturales."},
        {"num": "03", "header": "Compromiso",
         "text": "El lider semilla: una eleccion consciente."},
        {"num": "04", "header": "Mapa",
         "text": "El viaje de 7 modulos: del interior al sistema."},
    ])
    slides.append(("S2 Tarjetas", s))

    # ── S3: Separador Tema 01 ────────────────────────────────────────
    s = clone_slide(prs, T_SEPARATOR)
    build_separador(s, bg_sep, prs, "01",
                    "La tormenta perfecta",
                    "\u00BFQue esta pasando que hace liderar tan dificil?")
    slides.append(("S3 Separador 01", s))

    # ── S4: Visual central ───────────────────────────────────────────
    s = clone_slide(prs, T_VISUAL)
    build_visual_central(s, "APERTURA",
        "\u00BFQue significa liderar con consciencia?",
        "img-m00-00-lider-ante-tormenta.png",
        "Antes de empezar a trabajar sobre nosotros mismos, necesitamos "
        "mirar al contexto en el que estamos liderando. No como informacion "
        "abstracta, sino como un reconocimiento de lo que probablemente ya sientes.",
        "Este modulo es el portico: el espacio donde situamos el contexto "
        "antes de empezar el trabajo interior.")
    slides.append(("S4 Visual central", s))

    # ── S5: Grid 4 columnas ──────────────────────────────────────────
    s = clone_slide(prs, T_GRID)
    build_grid_cols(s, "DIAGNOSTICO", "Cuatro fuerzas convergentes", [
        {"header": "Fragmentacion digital",
         "text": "La atencion profunda ha desaparecido.",
         "stat": "47s de atencion"},
        {"header": "Crisis de legitimidad",
         "text": "Desconfianza estructural en el liderazgo.",
         "stat": "7 de cada 10"},
        {"header": "Irrupcion de la IA",
         "text": "Obsolescencia acelerada de habilidades.",
         "stat": "39% cambiaran"},
        {"header": "Epidemia de agotamiento",
         "text": "Estres cronico como norma laboral.",
         "stat": "82% en riesgo"},
    ], cita="Lo nuevo no es cada fuerza. Lo nuevo es su convergencia.")
    slides.append(("S5 Grid 4col", s))

    # ── S6: Stat 47 seg ──────────────────────────────────────────────
    s = clone_slide(prs, T_CARDS)
    build_stat(s, "47 seg",
        "de atencion sostenida en una pantalla",
        fuente="Gloria Mark, UC Irvine \u2014 Attention Span (2023)",
        apoyo="En 2004 eran 2.5 minutos. Hemos perdido dos tercios "
              "de nuestra capacidad de foco.",
        placeholder="Persona frente a pantallas con foco fragmentado")
    slides.append(("S6 Stat 47seg", s))

    # ── S7: Multi-stat fragmentacion ─────────────────────────────────
    s = clone_slide(prs, T_CARDS)
    build_multi_stat(s, "EVIDENCIA", "Fragmentacion digital", [
        {"num": "47 seg", "desc": "Promedio de atencion sostenida en pantalla.",
         "source": "Gloria Mark, UC Irvine (2023)"},
        {"num": "25 min", "desc": "Tiempo necesario para recuperar el foco tras una interrupcion.",
         "source": "Gloria Mark, UC Irvine (2023)"},
        {"num": "6x", "desc": "Velocidad a la que viajan las noticias falsas vs. la verdad.",
         "source": "MIT (Vosoughi et al., 2018)"},
        {"num": "39%", "desc": "De las competencias laborales cambiaran para 2030.",
         "source": "WEF Future of Jobs (2025)"},
    ])
    slides.append(("S7 Multi-stat", s))

    # ── S8: Stat 70% ─────────────────────────────────────────────────
    s = clone_slide(prs, T_CARDS)
    build_stat(s, "70%",
        "de la varianza en el compromiso de un equipo depende de su lider directo.",
        fuente="Gallup, State of the Global Workplace (2025)")
    slides.append(("S8 Stat 70%", s))

    # ── S9: Multi-stat agotamiento ───────────────────────────────────
    s = clone_slide(prs, T_CARDS)
    build_multi_stat(s, "EPIDEMIA", "Agotamiento global", [
        {"num": "82%", "desc": "de los lideres estan en riesgo de burnout.",
         "source": "Microsoft Work Trend Index"},
        {"num": "41%", "desc": "de los empleados reportan 'renuncia silenciosa'.",
         "source": "Gallup Global Report"},
        {"num": "58%", "desc": "siente estres diario en el trabajo.",
         "source": "State of the Global Workplace"},
    ], cita="OMS: 'El burnout ya no se clasifica como condicion medica, "
            "sino como fenomeno ocupacional resultado de estres cronico "
            "no gestionado.'")
    slides.append(("S9 Multi-stat agot", s))

    # ── S10: Causa-efecto ────────────────────────────────────────────
    s = clone_slide(prs, T_CAUSE)
    build_causa_efecto(s, "SISTEMA",
        "No son cuatro problemas. Es un sistema.",
        "Se retroalimentan mutuamente:",
        [
            "La fragmentacion reduce la capacidad de foco.",
            "Sin foco, aumenta la reactividad y el estres.",
            "La reactividad erosiona la calidad de la presencia.",
            "La mala presencia genera desconfianza en el equipo.",
            "La desconfianza aumenta la necesidad de control.",
        ],
        placeholder="ppt-m00-01-ciclo-retroalimentacion.png",
        cita="El ciclo no se rompe informandose mas. Se rompe estando mas presente.")
    slides.append(("S10 Causa-efecto", s))

    # ── S11: Cita Byung-Chul Han ────────────────────────────────────
    s = clone_slide(prs, T_CARDS)
    build_cita(s,
        "Ahora uno se explota a si mismo y cree que se esta realizando.",
        autor="Byung-Chul Han",
        obra="La sociedad del cansancio (2010)")
    slides.append(("S11 Cita", s))

    # ── S12: Separador Tema 02 ───────────────────────────────────────
    s = clone_slide(prs, T_SEPARATOR)
    build_separador(s, bg_sep, prs, "02",
                    "La tentacion de la impotencia",
                    "\u00BFQue hacemos con la sensacion de que no podemos hacer nada?")
    slides.append(("S12 Separador 02", s))

    # ── S13: Tarjetas 3 (respuestas naturales) ───────────────────────
    s = clone_slide(prs, T_CARDS)
    build_tarjetas(s, "RESPUESTAS NATURALES",
        "Cuando el problema parece mas grande que nosotros", [
        {"header": "Desconexion", "subtitle": "Esto no va conmigo.",
         "text": "Protege, pero aisla."},
        {"header": "Cinismo", "subtitle": "No tiene sentido.",
         "text": "Se siente como lucidez, pero paraliza."},
        {"header": "Agotamiento", "subtitle": "Si me esfuerzo mas, podre.",
         "text": "Parece compromiso, pero destruye."},
    ])
    slides.append(("S13 Tarjetas 3", s))

    # ── S14: Contraste dual ──────────────────────────────────────────
    s = clone_slide(prs, T_CONTRAST)
    build_contraste(s, "PATRON COMUN",
        "Nos separan de nuestros propios actos",
        "Quien eres",
        ["Tus valores", "Tu intencion", "Tu proposito"],
        "Como lideras",
        ["Tus decisiones", "Tus reacciones", "Tu impacto real"],
        cita="El problema no es que no sepamos liderar. "
             "Es que no estamos presentes cuando lo hacemos.",
        sym="\u2260")
    slides.append(("S14 Contraste", s))

    # ── S15: Reflexion ───────────────────────────────────────────────
    s = clone_slide(prs, T_SEPARATOR)
    build_reflexion(s, bg_sep, prs,
        "\u00BFCual de las tres respuestas reconoces mas en ti ultimamente?",
        apoyo="Una pausa para la honestidad.",
        opciones=["Desconexion", "Cinismo", "Agotamiento"],
        etiq="REFLEXION")
    slides.append(("S15 Reflexion", s))

    # ── S16: Separador Tema 03 ───────────────────────────────────────
    s = clone_slide(prs, T_SEPARATOR)
    build_separador(s, bg_sep, prs, "03",
                    "El lider semilla",
                    "\u00BFQue tipo de lider puedo ser en medio de esta tormenta?")
    slides.append(("S16 Separador 03", s))

    # ── S17: Proceso 4 fases ─────────────────────────────────────────
    s = clone_slide(prs, T_PROCESS)
    build_proceso(s, "METAFORA", "La naturaleza de la semilla", [
        {"num": 1, "header": "No elige el clima",
         "text": "Acepta las condiciones dadas sin resignacion."},
        {"num": 2, "header": "Echa raices",
         "text": "Trabaja en lo invisible antes de mostrar resultados."},
        {"num": 3, "header": "Transforma su entorno",
         "text": "Convierte lo que encuentra en nutriente."},
        {"num": 4, "header": "Crea condiciones",
         "text": "Su crecimiento facilita que otros tambien crezcan."},
    ], placeholder="img-m00-03-semilla-bajo-tierra.png")
    slides.append(("S17 Proceso", s))

    # ── S18: Contraste (lider semilla) ───────────────────────────────
    s = clone_slide(prs, T_CONTRAST)
    build_contraste(s, "DEFINICION",
        "Lo que define al lider semilla",
        "Lo que NO lo define",
        ["Conocimiento tecnico", "Autoridad formal",
         "Resultados visibles inmediatos"],
        "Lo que SI lo define",
        ["Calidad de su presencia",
         "Capacidad de no reaccionar automaticamente",
         "Poder sostener incertidumbre sin transmitir ansiedad",
         "Ver a las personas como personas, no como recursos"],
        cita="El lider semilla no elige entre ser buena persona y ser buen "
             "profesional. Ha descubierto que son la misma cosa.")
    slides.append(("S18 Contraste", s))

    # ── S19: Grid 5 items vertical ───────────────────────────────────
    s = clone_slide(prs, T_GRID)
    build_grid_rows(s, "NUEVAS CAPACIDADES", "Lo que la tormenta exige", [
        {"header": "Discernimiento", "text": "Ver lo que es, no lo que temes."},
        {"header": "Foco sostenido",
         "text": "Mantener la atencion en medio del ruido."},
        {"header": "Lectura del presente",
         "text": "Entender el contexto antes de actuar."},
        {"header": "Espacio interior",
         "text": "Crear una pausa entre estimulo y respuesta."},
        {"header": "Humanidad",
         "text": "Conectar con el otro como persona, no como recurso."},
    ], cita="No son habilidades blandas. Son habilidades criticas \u2014y entrenables.")
    slides.append(("S19 Grid rows", s))

    # ── S20: Cita variante (presencia) ───────────────────────────────
    s = clone_slide(prs, T_CARDS)
    build_cita(s,
        "La presencia no es solo etica. Es eficacia.",
        apoyo="Ve mas, porque no esta cegado por la prisa. Escucha mejor, "
              "porque no prepara su respuesta. Decide con precision, porque "
              "observa lo que pasa.",
        fuente="Liderar con Consciencia \u2014 El lider semilla")
    slides.append(("S20 Cita var", s))

    # ── S21: Reflexion (pregunta) ────────────────────────────────────
    s = clone_slide(prs, T_SEPARATOR)
    build_reflexion(s, bg_sep, prs,
        "\u00BFEstas dispuesto a trabajar en ti mismo?",
        apoyo="No hace falta una respuesta segura. Un si tentativo es suficiente.")
    slides.append(("S21 Reflexion", s))

    # ── S22: Separador Tema 04 ───────────────────────────────────────
    s = clone_slide(prs, T_SEPARATOR)
    build_separador(s, bg_sep, prs, "04",
                    "El mapa del viaje",
                    "\u00BFQue voy a aprender y como se conecta todo?")
    slides.append(("S22 Separador 04", s))

    # ── S23: Grafica + analisis ──────────────────────────────────────
    s = clone_slide(prs, T_CHART)
    build_grafica(s, "RECORRIDO", "Tres fases: del interior al sistema", [
        {"header": "Fase 1 (M1-M3)", "subtitle": "Consciencia Activada",
         "text": "El lider consigo mismo."},
        {"header": "Fase 2 (M4-M5)", "subtitle": "Consciencia Aplicada",
         "text": "El lider con el equipo."},
        {"header": "Fase 3 (M6-M7)", "subtitle": "Consciencia Dirigida",
         "text": "El lider con el sistema."},
    ], placeholder="graf-m00-04-tres-fases-v3.png")
    slides.append(("S23 Grafica", s))

    # ── S24: Roadmap 8 modulos ───────────────────────────────────────
    s = clone_slide(prs, T_ROADMAP)
    build_roadmap(s, "RECORRIDO", "El viaje de 7 modulos", [
        {"label": "M0", "subtitle": "El momento", "text": ""},
        {"label": "M1", "subtitle": "Lo que soy", "text": ""},
        {"label": "M2", "subtitle": "Como estoy", "text": ""},
        {"label": "M3", "subtitle": "Lo que pienso", "text": ""},
        {"label": "M4", "subtitle": "Lo que necesito", "text": ""},
        {"label": "M5", "subtitle": "Lo que necesitamos", "text": ""},
        {"label": "M6", "subtitle": "Lo que quiero", "text": ""},
        {"label": "M7", "subtitle": "Lo que esta bien", "text": ""},
    ])
    slides.append(("S24 Roadmap", s))

    # ── S25: Roadmap vertical (semilla) ──────────────────────────────
    s = clone_slide(prs, T_ROADMAP)
    build_roadmap(s, "ARCO DEL PROGRAMA", "La semilla y sus capacidades", [
        {"label": "M1", "subtitle": "Presencia",
         "text": "La semilla acepta su naturaleza y su lugar."},
        {"label": "M2", "subtitle": "Estabilidad",
         "text": "La semilla echa raices para sostenerse."},
        {"label": "M3", "subtitle": "Discernimiento",
         "text": "La semilla distingue la luz de la oscuridad."},
        {"label": "M4", "subtitle": "Sostenibilidad",
         "text": "La semilla absorbe los nutrientes que necesita."},
        {"label": "M5", "subtitle": "Conexion",
         "text": "La semilla crea un ecosistema con otras."},
        {"label": "M6", "subtitle": "Voluntad",
         "text": "La semilla crece con fuerza hacia su proposito."},
        {"label": "M7", "subtitle": "Integridad",
         "text": "La semilla da fruto y completa su ciclo."},
    ], cita="La promesa del lider semilla se cierra al final del programa.")
    slides.append(("S25 Roadmap vert", s))

    # ── S26: Tarjetas resumen (4) ────────────────────────────────────
    s = clone_slide(prs, T_CARDS)
    build_tarjetas(s, "RESUMEN Y CIERRE", "Ideas clave del modulo", [
        {"text": "El agotamiento actual es sistemico, no un fallo individual."},
        {"text": "Las defensas habituales (cinismo, desconexion) solo "
                 "agravan el aislamiento."},
        {"text": "El lider semilla transforma su entorno a traves de la "
                 "calidad de su presencia."},
        {"text": "No necesitamos certezas para empezar, solo la intencion "
                 "de trabajar en nosotros mismos."},
    ])
    slides.append(("S26 Tarjetas 4", s))

    # ── S27: Contraste De/A ──────────────────────────────────────────
    s = clone_slide(prs, T_CONTRAST)
    build_contraste(s, "TRANSFORMACION M0", "",
        "De:",
        ["La sensacion de impotencia ante la magnitud de la tormenta."],
        "A:",
        ["La certeza de que nuestra presencia es la primera "
         "herramienta de cambio."])
    slides.append(("S27 Contraste De/A", s))

    # ── S28: Separador practica ──────────────────────────────────────
    s = clone_slide(prs, T_PRACTICE)
    build_sep_practica(s, bg_pract, prs,
        "Meditacion de\nintencion",
        "Conectar con el proposito de estar aqui.",
        detalle="~ 5 minutos")
    slides.append(("S28 Sep practica", s))

    # ── S29: Transicion ──────────────────────────────────────────────
    s = clone_slide(prs, T_TRANSITION)
    build_transicion(s, bg_trans, prs,
        "MODULO 01", "Consciente de lo que Soy",
        "Las semillas no necesitan certezas para germinar. "
        "Solo necesitan empezar.",
        placeholder="img-m00-00-lider-semilla-brote.png")
    slides.append(("S29 Transicion", s))

    # ── S30: Cierre ──────────────────────────────────────────────────
    s = clone_slide(prs, T_SEPARATOR)
    build_cierre(s, bg_sep, prs)
    slides.append(("S30 Cierre", s))

    # ── Remove original template slides ──────────────────────────────
    print(f"Removing {num_template} original template slides...")
    remove_slides(prs, list(range(num_template)))

    # ── Transfer speaker notes from M00 v1 ───────────────────────────
    print("Transferring speaker notes...")
    notes_count = 0
    for i, slide in enumerate(prs.slides):
        if i < len(m00v1.slides):
            src = m00v1.slides[i]
            if src.has_notes_slide:
                src_tf = src.notes_slide.notes_text_frame
                if src_tf is None:
                    continue
                src_text = src_tf.text
                if not src_text or not src_text.strip():
                    continue
                try:
                    tgt_ns = slide.notes_slide
                    tgt_tf = tgt_ns.notes_text_frame
                    if tgt_tf is not None:
                        tgt_tf.text = src_text
                        notes_count += 1
                    else:
                        # Fallback: inject text into notes XML body directly
                        nsmap = {'a': NS_A, 'p': NS_P}
                        body = tgt_ns._element.find('.//p:txBody', nsmap)
                        if body is None:
                            # Find any txBody in the notes
                            body = tgt_ns._element.find(f'.//{{{NS_P}}}sp/{{{NS_P}}}txBody')
                        if body is None:
                            # Create a body placeholder in the notes
                            cSld = tgt_ns._element.find(f'{{{NS_P}}}cSld')
                            if cSld is None:
                                continue
                            spTree = cSld.find(f'{{{NS_P}}}spTree')
                            if spTree is None:
                                continue
                            # Add a simple sp with txBody
                            sp = etree.SubElement(spTree, f'{{{NS_P}}}sp')
                            nvSpPr = etree.SubElement(sp, f'{{{NS_P}}}nvSpPr')
                            cNvPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvPr',
                                                     attrib={'id': '3', 'name': 'Notes Placeholder'})
                            cNvSpPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvSpPr')
                            nvPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}nvPr')
                            ph = etree.SubElement(nvPr, f'{{{NS_P}}}ph',
                                                  attrib={'type': 'body', 'idx': '1'})
                            spPr = etree.SubElement(sp, f'{{{NS_P}}}spPr')
                            body = etree.SubElement(sp, f'{{{NS_P}}}txBody')
                            bodyPr = etree.SubElement(body, f'{{{NS_A}}}bodyPr')
                            lstStyle = etree.SubElement(body, f'{{{NS_A}}}lstStyle')

                        # Clear existing paragraphs and add text
                        for old_p in list(body.findall(f'{{{NS_A}}}p')):
                            body.remove(old_p)
                        for line in src_text.split('\n'):
                            p_el = etree.SubElement(body, f'{{{NS_A}}}p')
                            r_el = etree.SubElement(p_el, f'{{{NS_A}}}r')
                            t_el = etree.SubElement(r_el, f'{{{NS_A}}}t')
                            t_el.text = line
                        notes_count += 1
                except Exception as e:
                    print(f"  Warning: Could not set notes for slide {i+1}: {e}")
    print(f"  Transferred {notes_count} speaker notes")

    # ── Save ─────────────────────────────────────────────────────────
    OUTPUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_FILE))
    print(f"\nSaved: {OUTPUT_FILE}")

    # ── Verify ───────────────────────────────────────────────────────
    return verify()


# ═══════════════════════════════════════════════════════════════════════════
# SECTION 5: VERIFICATION
# ═══════════════════════════════════════════════════════════════════════════

def verify():
    """Verify the output PPTX."""
    print("\n=== VERIFICATION ===")
    prs = Presentation(str(OUTPUT_FILE))
    ok = True

    # 1. Slide count
    n = len(prs.slides)
    status = "OK" if n == 30 else "FAIL"
    if n != 30:
        ok = False
    print(f"  Slide count: {n} [{status}]")

    # 2. Dimensions
    dims_ok = prs.slide_width == SLIDE_W and prs.slide_height == SLIDE_H
    status = "OK" if dims_ok else "FAIL"
    if not dims_ok:
        ok = False
    print(f"  Dimensions: {prs.slide_width}x{prs.slide_height} [{status}]")

    # 3. Check for M04 residues
    m04_found = False
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "Consciente de lo" in text and "Necesito" in text:
                    print(f"  WARNING: M04 residue in slide {i+1}: {text[:60]}")
                    m04_found = True
    status = "OK" if not m04_found else "WARN"
    print(f"  No M04 residue: [{status}]")

    # 4. Speaker notes
    def has_notes(s):
        try:
            if not s.has_notes_slide:
                return False
            tf = s.notes_slide.notes_text_frame
            return tf is not None and tf.text.strip() != ""
        except Exception:
            return False

    notes_count = sum(1 for s in prs.slides if has_notes(s))
    expected = 29  # All except S30
    status = "OK" if notes_count >= expected else f"WARN ({notes_count}/{expected})"
    print(f"  Speaker notes: {notes_count} [{status}]")

    # 5. Print slide summary
    print("\n  Slide summary:")
    for i, slide in enumerate(prs.slides):
        first_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t and len(t) > 2:
                    first_text = t[:50]
                    break
        n_flag = "N" if has_notes(slide) else "-"
        print(f"    S{i+1:2d} [{n_flag}] {first_text}")

    if ok:
        print("\n  ALL CHECKS PASSED")
    return ok


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    success = build_m00()
    sys.exit(0 if success else 1)
