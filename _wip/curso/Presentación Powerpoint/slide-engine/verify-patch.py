#!/usr/bin/env python
"""
Verify PPTX Patch — 6 automated checks per patched PPTX.

Usage:
    python verify-patch.py M06          # Verify a single module
    python verify-patch.py --all        # Verify all patched modules

Checks:
    1. SLIDE_COUNT   - Patched slide count >= original
    2. PORTADA_TEXT   - Slide 1 contains correct module tag
    3. GRACIAS_EXISTS - Penultimate or last slide contains "GRACIAS"
    4. REFS_EXISTS    - Last slide contains "PROFUNDIZAR" or "Fuentes" (if refs expected)
    5. CONTENT_PRESERVED - Non-structural slides have same text as original
    6. BG_IMAGES     - Structural slides have background images present
"""

import sys
import json
import hashlib
from pathlib import Path
from pptx import Presentation

# ── Paths ──────────────────────────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).parent
MANUS_DIR = SCRIPT_DIR.parent / "curso" / "Presentación Powerpoint" / "Diapos Manus"
OUTPUT_DIR = SCRIPT_DIR.parent / "curso" / "Presentación Powerpoint" / "ppts CLAUDE CODE"
MODULE_DATA_FILE = SCRIPT_DIR / "module-data.json"


def load_module_data():
    with open(MODULE_DATA_FILE, "r", encoding="utf-8") as f:
        return json.load(f)


def get_slide_texts(slide):
    """Extract all text from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def get_structural_indices(data):
    """Get set of 0-based indices that are structural slides."""
    indices = set()
    indices.add(0)  # Portada is always slide 0

    for sep in data.get("separators", []):
        indices.add(sep["slide_index"])

    gracias = data.get("gracias")
    if gracias and gracias.get("slide_index") is not None:
        indices.add(gracias["slide_index"])

    refs = data.get("refs")
    if refs and refs.get("slide_index") is not None:
        indices.add(refs["slide_index"])

    return indices


def check_slide_count(original_prs, patched_prs, data):
    """Check 1: Patched slide count >= original."""
    orig_count = len(original_prs.slides)
    patched_count = len(patched_prs.slides)
    passed = patched_count >= orig_count
    detail = f"original={orig_count}, patched={patched_count}"
    if not passed:
        detail += f" (MISSING {orig_count - patched_count} slides!)"
    return passed, detail


def check_portada_text(patched_prs, data):
    """Check 2: Slide 1 contains the correct module tag."""
    slide = patched_prs.slides[0]
    texts = get_slide_texts(slide)
    tag = data["portada"]["tag"]
    found = any(tag in t for t in texts)
    detail = f"expected tag='{tag}', found in texts={found}"
    if not found:
        detail += f", actual texts: {texts[:3]}"
    return found, detail


def check_gracias_exists(patched_prs, data):
    """Check 3: Penultimate or last slide contains 'GRACIAS'."""
    n = len(patched_prs.slides)
    for i in range(max(0, n - 3), n):
        texts = get_slide_texts(patched_prs.slides[i])
        text_joined = " ".join(texts).upper()
        if "GRACIAS" in text_joined:
            return True, f"Found 'GRACIAS' in slide {i + 1}"
    return False, f"'GRACIAS' not found in last 3 slides"


def check_refs_exists(patched_prs, data):
    """Check 4: Last slide contains refs keywords (only if refs expected)."""
    refs = data.get("refs")
    if not refs or refs.get("slide_index") is None:
        return True, "No refs expected for this module (SKIP)"

    n = len(patched_prs.slides)
    for i in range(max(0, n - 2), n):
        texts = get_slide_texts(patched_prs.slides[i])
        text_joined = " ".join(texts).upper()
        if "PROFUNDIZAR" in text_joined or "FUENTES" in text_joined:
            return True, f"Found refs keywords in slide {i + 1}"
    return False, "Refs keywords not found in last 2 slides"


def check_content_preserved(original_prs, patched_prs, data):
    """Check 5: Non-structural slides have same text as original."""
    structural = get_structural_indices(data)
    orig_slides = list(original_prs.slides)
    patched_slides = list(patched_prs.slides)

    # Content slides are all non-structural slides in the original
    mismatches = []
    for i in range(len(orig_slides)):
        if i in structural:
            continue
        if i >= len(patched_slides):
            mismatches.append(f"slide {i + 1}: missing in patched")
            continue

        orig_texts = get_slide_texts(orig_slides[i])
        patched_texts = get_slide_texts(patched_slides[i])

        # Compare text content (join and normalize)
        orig_str = " ".join(orig_texts).strip()
        patched_str = " ".join(patched_texts).strip()

        if orig_str != patched_str:
            mismatches.append(f"slide {i + 1}")

    if not mismatches:
        content_count = len(orig_slides) - len(structural)
        return True, f"All {content_count} content slides preserved"
    else:
        return False, f"Mismatches in: {', '.join(mismatches[:5])}"


def check_bg_images(patched_prs, data):
    """Check 6: Structural slides have background images."""
    structural = get_structural_indices(data)
    # Also check added slides (gracias/refs at end)
    n = len(patched_prs.slides)

    missing_bg = []
    for i in range(n):
        # Check structural slides and last 2 slides
        if i in structural or i >= n - 2:
            slide = patched_prs.slides[i]
            has_image = False
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    has_image = True
                    break
            if not has_image:
                # Also check rels for images
                for rel in slide.part.rels.values():
                    if "image" in str(rel.reltype):
                        has_image = True
                        break
            if not has_image:
                missing_bg.append(f"slide {i + 1}")

    if not missing_bg:
        return True, f"All structural slides have background images"
    else:
        return False, f"Missing bg in: {', '.join(missing_bg)}"


def verify_module(module_id, module_data):
    """Run all 6 checks on a patched module."""
    data = module_data[module_id]
    if data.get("output_filename") is None:
        print(f"  Skipping {module_id} (reference module)")
        return None

    original_path = MANUS_DIR / data["filename"]
    patched_path = OUTPUT_DIR / data["output_filename"]

    if not original_path.exists():
        print(f"  ERROR: Original not found: {original_path}")
        return False
    if not patched_path.exists():
        print(f"  ERROR: Patched not found: {patched_path}")
        return False

    original_prs = Presentation(str(original_path))
    patched_prs = Presentation(str(patched_path))

    checks = [
        ("SLIDE_COUNT", lambda: check_slide_count(original_prs, patched_prs, data)),
        ("PORTADA_TEXT", lambda: check_portada_text(patched_prs, data)),
        ("GRACIAS_EXISTS", lambda: check_gracias_exists(patched_prs, data)),
        ("REFS_EXISTS", lambda: check_refs_exists(patched_prs, data)),
        ("CONTENT_PRESERVED", lambda: check_content_preserved(original_prs, patched_prs, data)),
        ("BG_IMAGES", lambda: check_bg_images(patched_prs, data)),
    ]

    passed = 0
    total = len(checks)
    results = []

    for name, check_fn in checks:
        try:
            ok, detail = check_fn()
            status = "PASS" if ok else "FAIL"
            if ok:
                passed += 1
            results.append((name, status, detail))
            print(f"  [{status}] {name}: {detail}")
        except Exception as e:
            results.append((name, "ERROR", str(e)))
            print(f"  [ERROR] {name}: {e}")

    verdict = f"{passed}/{total}"
    all_passed = passed == total
    print(f"\n  Result: {'PASS' if all_passed else 'FAIL'} {verdict}")
    return all_passed


def main():
    if len(sys.argv) < 2:
        print("Usage: python verify-patch.py <MODULE_ID> | --all")
        sys.exit(1)

    module_data = load_module_data()

    if sys.argv[1] == "--all":
        order = ["M06", "M05", "M04", "M03", "M02", "M01D", "M01C", "M01B", "M01A", "M00"]
        results = {}
        for mod in order:
            if mod in module_data and module_data[mod].get("output_filename"):
                patched_path = OUTPUT_DIR / module_data[mod]["output_filename"]
                if patched_path.exists():
                    print(f"\n{'='*60}")
                    print(f"Verifying {mod}...")
                    print(f"{'='*60}")
                    results[mod] = verify_module(mod, module_data)

        print(f"\n{'='*60}")
        print("SUMMARY")
        print(f"{'='*60}")
        for mod, ok in results.items():
            status = "PASS" if ok else "FAIL"
            print(f"  {mod}: {status}")
    else:
        module_id = sys.argv[1].upper()
        if module_id not in module_data:
            print(f"ERROR: Unknown module '{module_id}'")
            sys.exit(1)
        print(f"Verifying {module_id}...")
        ok = verify_module(module_id, module_data)
        sys.exit(0 if ok else 1)


if __name__ == "__main__":
    main()
