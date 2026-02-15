#!/usr/bin/env python
"""
Merge M01 sub-blocks (M01A, M01B, M01C, M01D) into a single PPTX.

- M01A portada is re-patched to "Módulo 01" / "Consciente de\nlo que Soy"
- Intermediate Gracias slides are removed
- M01D's Gracias is kept as the final Gracias
"""

import copy
import sys
from pathlib import Path
from io import BytesIO
from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

SCRIPT_DIR = Path(__file__).parent
OUTPUT_DIR = SCRIPT_DIR.parent / "curso" / "Presentación Powerpoint" / "ppts CLAUDE CODE"

# Use the patched PPTXs as source
FILES = {
    "M01A": OUTPUT_DIR / "M01A-El-Cerebro-Constructor.pptx",
    "M01B": OUTPUT_DIR / "M01B-El-Piloto-Automatico.pptx",
    "M01C": OUTPUT_DIR / "M01C-La-Puerta-de-Salida.pptx",
    "M01D": OUTPUT_DIR / "M01D-El-Viaje-y-el-Compromiso.pptx",
}

OUTPUT_FILE = OUTPUT_DIR / "M01-Consciente-de-lo-que-Soy.pptx"

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
REL_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
REL_SLIDE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"


def get_slide_texts(slide):
    """Extract all text from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def is_gracias_slide(slide):
    """Check if a slide is a Gracias slide."""
    texts = get_slide_texts(slide)
    return any("gracias" in t.lower() for t in texts)


def copy_slide(src_prs, src_slide, dst_prs):
    """Copy a slide from src_prs to dst_prs using XML manipulation."""
    # Get a blank slide layout from destination
    slide_layout = None
    for layout in dst_prs.slide_layouts:
        if layout.name.lower() in ("blank", "en blanco"):
            slide_layout = layout
            break
    if slide_layout is None:
        slide_layout = dst_prs.slide_layouts[-1]

    new_slide = dst_prs.slides.add_slide(slide_layout)

    # Clear default shapes from new slide
    spTree = new_slide.shapes._spTree
    for child in list(spTree):
        tag = etree.QName(child).localname
        if tag in ("sp", "pic", "grpSp", "cxnSp"):
            spTree.remove(child)

    # Copy image relationships first — build a mapping of old rId to new rId
    rId_map = {}
    for rel in src_slide.part.rels.values():
        if "image" in str(rel.reltype):
            # Copy the image blob to destination
            image_blob = rel.target_part.blob
            content_type = rel.target_part.content_type

            # Add image to new slide's package
            img_stream = BytesIO(image_blob)
            # Use the slide part to add the image relationship
            image_part, new_rId = new_slide.part.get_or_add_image_part(img_stream)
            rId_map[rel.rId] = new_rId

    # Copy all shape elements from source to destination
    src_spTree = src_slide.shapes._spTree
    for child in list(src_spTree):
        tag = etree.QName(child).localname
        if tag in ("sp", "pic", "grpSp", "cxnSp"):
            new_element = copy.deepcopy(child)

            # Update image rIds in copied elements
            for blip in new_element.findall(f".//{{{NS_A}}}blip"):
                old_rId = blip.get(f"{{{NS_R}}}embed")
                if old_rId and old_rId in rId_map:
                    blip.set(f"{{{NS_R}}}embed", rId_map[old_rId])

            spTree.append(new_element)

    # Copy slide background if present
    src_bg = src_slide._element.find(f"{{{NS_P}}}bg")
    if src_bg is not None:
        new_bg = copy.deepcopy(src_bg)
        # Insert bg before spTree
        new_slide._element.insert(0, new_bg)

    return new_slide


def repatch_portada(slide, tag_text, title_text, quote_text):
    """Update text in the portada slide (first 3 text shapes)."""
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_shapes.append(shape)

    # Sort by top position to identify tag, title, quote
    text_shapes.sort(key=lambda s: s.top)

    if len(text_shapes) >= 3:
        # Shape 0 = tag, Shape 1 = title, Shape 2 = quote
        for i, (shape, new_text) in enumerate(zip(text_shapes[:3], [tag_text, title_text, quote_text])):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
            # Set text in first paragraph's first run
            lines = new_text.split("\n")
            for j, line in enumerate(lines):
                if j == 0:
                    if shape.text_frame.paragraphs[0].runs:
                        shape.text_frame.paragraphs[0].runs[0].text = line
                    else:
                        run = shape.text_frame.paragraphs[0].add_run()
                        run.text = line
                else:
                    p = shape.text_frame.add_paragraph()
                    # Copy formatting from first paragraph
                    src_para = shape.text_frame.paragraphs[0]
                    p.alignment = src_para.alignment
                    run = p.add_run()
                    run.text = line
                    if src_para.runs:
                        src_run = src_para.runs[0]
                        run.font.name = src_run.font.name
                        run.font.size = src_run.font.size
                        run.font.bold = src_run.font.bold
                        if src_run.font.color and src_run.font.color.rgb:
                            run.font.color.rgb = src_run.font.color.rgb
        print(f"  Portada re-patched: tag='{tag_text}', title='{title_text}'")
    else:
        print(f"  WARNING: Only {len(text_shapes)} text shapes found in portada")


def main():
    # Verify all input files exist
    for mod, path in FILES.items():
        if not path.exists():
            print(f"ERROR: {mod} not found at {path}")
            sys.exit(1)
        print(f"  Found {mod}: {path.name}")

    # Load all presentations
    print("\nLoading presentations...")
    prss = {}
    for mod, path in FILES.items():
        prss[mod] = Presentation(str(path))
        print(f"  {mod}: {len(prss[mod].slides)} slides")

    # Use M01A as the base
    base_prs = prss["M01A"]
    base_slides = list(base_prs.slides)
    print(f"\nBase: M01A with {len(base_slides)} slides")

    # Identify and remove Gracias slide from M01A (last slide)
    if is_gracias_slide(base_slides[-1]):
        print("  Removing M01A Gracias slide (last slide)")
        # Remove last slide from presentation XML
        rId = base_prs.slides._sldIdLst[-1].get(f"{{{NS_R}}}id")
        sldId = base_prs.slides._sldIdLst[-1]
        base_prs.slides._sldIdLst.remove(sldId)
        # Also remove the relationship
        try:
            base_prs.part.drop_rel(rId)
        except (KeyError, ValueError):
            pass
        print(f"  M01A now has {len(base_prs.slides)} slides")
    else:
        print("  WARNING: Last slide of M01A is NOT Gracias")

    # Re-patch M01A portada
    print("\nRe-patching M01A portada...")
    repatch_portada(
        base_prs.slides[0],
        "Módulo 01",
        "Consciente de\nlo que Soy",
        "\"El cerebro que crees ser vs. el cerebro que realmente eres.\""
    )

    # Copy slides from M01B, M01C, M01D
    for mod in ["M01B", "M01C", "M01D"]:
        src_prs = prss[mod]
        src_slides = list(src_prs.slides)
        print(f"\nCopying {mod} ({len(src_slides)} slides)...")

        for i, slide in enumerate(src_slides):
            # Skip Gracias slides for M01B and M01C
            if mod != "M01D" and is_gracias_slide(slide):
                print(f"  Skipping {mod} slide {i+1} (Gracias)")
                continue

            copy_slide(src_prs, slide, base_prs)
            texts = get_slide_texts(slide)
            label = texts[0][:40] if texts else "(empty)"
            print(f"  Copied {mod} slide {i+1}: {label}")

    # Save
    print(f"\nTotal slides in merged PPTX: {len(base_prs.slides)}")
    OUTPUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    base_prs.save(str(OUTPUT_FILE))
    print(f"Saved: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
