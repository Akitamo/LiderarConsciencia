#!/usr/bin/env python
"""
PPTX Patcher — Replaces structural slides in Manus PPTX files with M07-style design.

Usage:
    python pptx-patcher.py M06          # Patch a single module
    python pptx-patcher.py --all        # Patch all modules (including M07 theme-only)

Rules:
    - Text comes ONLY from module-data.json (extracted from original PPTX)
    - M07 is the visual reference (backgrounds + layout)
    - Content slides are NEVER modified
"""

import sys
import os
import json
import copy
import hashlib
import re
from io import BytesIO
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).parent
MANUS_DIR = SCRIPT_DIR.parent / "curso" / "Presentación Powerpoint" / "Diapos Manus"
OUTPUT_DIR = SCRIPT_DIR.parent / "curso" / "Presentación Powerpoint" / "ppts CLAUDE CODE"
M07_FILE = MANUS_DIR / "m07-ppt-Consciente_de_lo_que_Está_Bien.pptx"
MODULE_DATA_FILE = SCRIPT_DIR / "module-data.json"

# ── M07 Template Constants ────────────────────────────────────────────────
SLIDE_W = 9144000
SLIDE_H = 5143500

# Colors (from M07)
COLOR_ACCENT  = RGBColor(0xDC, 0x80, 0x60)
COLOR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SUBTLE  = RGBColor(0xE0, 0xEB, 0xEB)
COLOR_META    = RGBColor(0xA8, 0xB0, 0xA8)
COLOR_REF_HDR = RGBColor(0x2A, 0x50, 0x58)
COLOR_REF_BDY = RGBColor(0x5A, 0x6A, 0x6A)
COLOR_REF_CPT = RGBColor(0x88, 0x90, 0x88)

FONT_NAME = "Montserrat"

# XML namespaces
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
REL_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"


def load_module_data():
    with open(MODULE_DATA_FILE, "r", encoding="utf-8") as f:
        return json.load(f)


def extract_m07_backgrounds():
    """Extract background JPEG bytes from M07 structural slides."""
    prs = Presentation(str(M07_FILE))
    backgrounds = {}
    structural_indices = {
        "portada": 0,
        "separator": 6,
        "gracias": 41,
        "refs": 42,
    }
    for stype, idx in structural_indices.items():
        slide = prs.slides[idx]
        for rel in slide.part.rels.values():
            if "image" in str(rel.reltype):
                backgrounds[stype] = rel.target_part.blob
                break
    return backgrounds


def clear_slide_shapes(slide):
    """Remove all shapes from a slide AND clean up orphaned image relationships."""
    spTree = slide.shapes._spTree

    # 1. Collect rIds of images being removed
    orphan_rIds = set()
    for child in list(spTree):
        tag = etree.QName(child).localname
        if tag == "pic":
            blip = child.find(f".//{{{NS_A}}}blip")
            if blip is not None:
                rId = blip.get(f"{{{NS_R}}}embed")
                if rId:
                    orphan_rIds.add(rId)

    # 2. Remove shape elements
    for child in list(spTree):
        tag = etree.QName(child).localname
        if tag in ("sp", "pic", "grpSp", "cxnSp"):
            spTree.remove(child)

    # 3. Remove orphaned image relationships from slide rels
    for rId in orphan_rIds:
        try:
            slide.part.drop_rel(rId)
        except (KeyError, ValueError):
            pass


def add_bg_image(slide, jpeg_bytes, prs):
    """Add a full-slide background JPEG image."""
    img_stream = BytesIO(jpeg_bytes)
    pic = slide.shapes.add_picture(img_stream, 0, 0, prs.slide_width, prs.slide_height)
    # Move the picture to the back
    spTree = slide.shapes._spTree
    sp_element = pic._element
    spTree.remove(sp_element)
    spTree.insert(2, sp_element)
    return pic


def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color=COLOR_WHITE, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Emu(font_size)
        run.font.bold = bold
        run.font.color.rgb = color

    return txBox


def patch_theme_fonts(prs, font_name="Montserrat"):
    """Replace theme major/minor fonts with Montserrat."""
    for rel in prs.part.rels.values():
        if "theme" in str(rel.reltype):
            theme_part = rel.target_part
            xml_bytes = theme_part.blob
            xml_str = xml_bytes.decode("utf-8")

            # Replace major font (headings): <a:latin typeface="Calibri Light"/>
            xml_str = re.sub(
                r'(<a:majorFont[^>]*>.*?<a:latin\s+typeface=")[^"]+(")',
                rf'\g<1>{font_name}\2',
                xml_str,
                flags=re.DOTALL
            )
            # Replace minor font (body): <a:latin typeface="Calibri"/>
            xml_str = re.sub(
                r'(<a:minorFont[^>]*>.*?<a:latin\s+typeface=")[^"]+(")',
                rf'\g<1>{font_name}\2',
                xml_str,
                flags=re.DOTALL
            )

            theme_part._blob = xml_str.encode("utf-8")
            return True
    return False


def build_portada(slide, prs, bg_bytes, data):
    """Build a M07-style portada slide."""
    clear_slide_shapes(slide)
    add_bg_image(slide, bg_bytes, prs)

    add_textbox(slide, 714375, 1006041, 6786563, 233958,
                data["tag"], 151511, bold=True, color=COLOR_ACCENT)
    add_textbox(slide, 714375, 1454311, 6786563, 1200150,
                data["title"], 526415, bold=True, color=COLOR_WHITE)
    add_textbox(slide, 1035844, 3497424, 6465094, 640035,
                data["quote"], 216408, bold=False, color=COLOR_SUBTLE)


def build_separator(slide, prs, bg_bytes, number, title, subtitle):
    """Build a M07-style separator slide."""
    clear_slide_shapes(slide)
    add_bg_image(slide, bg_bytes, prs)

    add_textbox(slide, 714375, 1617166, 1471054, 1028700,
                number, 1184529, bold=True, color=COLOR_ACCENT)
    add_textbox(slide, 3042679, 1688604, 4920816, 194667,
                "TEMA", 125349, bold=True, color=COLOR_SUBTLE)
    add_textbox(slide, 3042679, 2026146, 4920816, 1005818,
                title, 418338, bold=True, color=COLOR_WHITE)
    add_textbox(slide, 3042679, 3246276, 4920816, 280029,
                subtitle, 188722, bold=False, color=COLOR_SUBTLE)


def build_gracias(slide, prs, bg_bytes):
    """Build a M07-style gracias slide."""
    clear_slide_shapes(slide)
    add_bg_image(slide, bg_bytes, prs)

    add_textbox(slide, 2594883, 1456432, 3954205, 857250,
                "Gracias", 789686, bold=True, color=COLOR_ACCENT,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 2594883, 2770882, 3954205, 312539,
                "Liderar con Consciencia", 216408, bold=False, color=COLOR_WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 2594883, 3512046, 3954205, 175022,
                "www.liderarconconsciencia.com", 119634, bold=False, color=COLOR_META,
                alignment=PP_ALIGN.CENTER)


def build_refs(slide, prs, bg_bytes, sources):
    """Build a M07-style refs slide with 2-column layout."""
    clear_slide_shapes(slide)
    add_bg_image(slide, bg_bytes, prs)

    # Header
    add_textbox(slide, 571500, 428625, 8001000, 155377,
                "PARA PROFUNDIZAR", 99568, bold=True, color=COLOR_ACCENT)
    add_textbox(slide, 571500, 655439, 8001000, 408980,
                "Fuentes principales", 269367, bold=True, color=COLOR_REF_HDR)

    # Divider line
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(571500), Emu(1207294), Emu(8001000), Emu(14288)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_REF_HDR
    shape.line.fill.background()

    # Sources in 2-column layout
    col_left_x = 735806
    col_right_x = 4950619
    col_width = 3621881
    start_y = 1564481
    row_height = 435844

    for i, source in enumerate(sources):
        col = i % 2
        row = i // 2
        x = col_left_x if col == 0 else col_right_x
        y = start_y + row * row_height
        add_textbox(slide, x, y, col_width, 300000,
                    source, 110000, bold=False, color=COLOR_REF_BDY)


def add_blank_slide(prs):
    """Add a blank slide at the end of the presentation."""
    slide_layout = None
    for layout in prs.slide_layouts:
        if layout.name.lower() in ("blank", "en blanco", "vacío", "vacio"):
            slide_layout = layout
            break
    if slide_layout is None:
        slide_layout = prs.slide_layouts[-1]

    slide = prs.slides.add_slide(slide_layout)
    clear_slide_shapes(slide)
    return slide


def patch_module(module_id, module_data, backgrounds):
    """Patch a single module's PPTX."""
    data = module_data[module_id]

    input_path = MANUS_DIR / data["filename"]
    output_filename = data.get("output_filename")

    # M07: only patch theme, no structural changes
    if output_filename is None:
        output_filename = "M07-Consciente-de-lo-que-Esta-Bien.pptx"

    output_path = OUTPUT_DIR / output_filename

    if not input_path.exists():
        print(f"  ERROR: Input file not found: {input_path}")
        return False

    print(f"  Input:  {input_path}")
    print(f"  Output: {output_path}")

    prs = Presentation(str(input_path))
    slides = prs.slides
    total = len(slides)
    patched_slides = []

    # Always patch theme fonts to Montserrat
    if patch_theme_fonts(prs, FONT_NAME):
        print(f"  Theme fonts patched to {FONT_NAME}")
        patched_slides.append("Theme")

    # M07: no structural patching needed
    if module_id == "M07":
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        print(f"  Saved: {output_path}")
        print(f"  Patched: {', '.join(patched_slides)}")
        return True

    # 1. Patch portada (slide 0)
    print(f"  Patching portada (slide 1)...")
    build_portada(slides[0], prs, backgrounds["portada"], data["portada"])
    patched_slides.append("Portada")

    # 2. Patch separators
    for sep in data.get("separators", []):
        idx = sep["slide_index"]
        if idx < total:
            print(f"  Patching separator (slide {idx + 1}): {sep['number']} - {sep['title'][:30]}...")
            build_separator(slides[idx], prs, backgrounds["separator"],
                           sep["number"], sep["title"], sep["subtitle"])
            patched_slides.append(f"Sep {sep['number']}")
        else:
            print(f"  WARNING: Separator slide index {idx} out of range (total={total})")

    # 3. Handle gracias
    gracias_info = data.get("gracias")
    if gracias_info and gracias_info.get("slide_index") is not None:
        idx = gracias_info["slide_index"]
        if idx < total:
            print(f"  Patching gracias (slide {idx + 1})...")
            build_gracias(slides[idx], prs, backgrounds["gracias"])
            patched_slides.append("Gracias")
        else:
            print(f"  WARNING: Gracias slide index {idx} out of range")
    else:
        print(f"  Adding gracias slide at the end...")
        new_slide = add_blank_slide(prs)
        build_gracias(new_slide, prs, backgrounds["gracias"])
        patched_slides.append("Gracias (added)")

    # 4. Handle refs
    refs_info = data.get("refs")
    if refs_info and refs_info.get("slide_index") is not None:
        idx = refs_info["slide_index"]
        sources = refs_info.get("sources", [])
        if idx < total and sources:
            print(f"  Patching refs (slide {idx + 1})...")
            build_refs(slides[idx], prs, backgrounds["refs"], sources)
            patched_slides.append("Refs")
        else:
            print(f"  WARNING: Refs slide index {idx} out of range or no sources")

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"  Saved: {output_path}")
    print(f"  Patched: {', '.join(patched_slides)}")
    return True


def main():
    if len(sys.argv) < 2:
        print("Usage: python pptx-patcher.py <MODULE_ID> | --all")
        print("Example: python pptx-patcher.py M06")
        sys.exit(1)

    module_data = load_module_data()
    print("Extracting M07 backgrounds...")
    backgrounds = extract_m07_backgrounds()
    print(f"  Extracted {len(backgrounds)} backgrounds: {list(backgrounds.keys())}")

    if sys.argv[1] == "--all":
        order = ["M07", "M06", "M05", "M04", "M03", "M02", "M01D", "M01C", "M01B", "M01A", "M00"]
        for mod in order:
            if mod in module_data:
                print(f"\n{'='*60}")
                print(f"Patching {mod}...")
                print(f"{'='*60}")
                patch_module(mod, module_data, backgrounds)
    else:
        module_id = sys.argv[1].upper()
        if module_id not in module_data:
            print(f"ERROR: Unknown module '{module_id}'. Available: {list(module_data.keys())}")
            sys.exit(1)
        print(f"\nPatching {module_id}...")
        patch_module(module_id, module_data, backgrounds)


if __name__ == "__main__":
    main()
