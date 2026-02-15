"""
Extract structured slide descriptions from Manus PPTX files.

Generates a markdown "visual script" for each presentation, capturing
the hierarchy of each slide: titles, subtitles, bullet points,
multi-column layouts, stats, quotes, and image placeholders.

Usage:
    python extract-slide-descriptions.py M02          # single module
    python extract-slide-descriptions.py --all        # all 6 modules (M02-M07)
    python extract-slide-descriptions.py M02 M05 M07  # specific modules
"""
import argparse
import os
import re
import sys
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
EMU_PER_INCH = 914400
SLIDE_W_EMU = 9144000   # 10"
SLIDE_H_EMU = 5143500   # 5.625"
BG_TOLERANCE = 50000    # ~0.05"

# Slide-bottom zone: shapes below this Y threshold with small font and
# short numeric text are treated as decorative slide numbers.
BOTTOM_ZONE_Y = SLIDE_H_EMU * 0.88  # bottom 12% of slide

# Font-size thresholds (in points) for role classification.
# These were determined empirically from the Manus PPTX files.
FONT_SECTION_NUMBER = 39   # ≥ this and short numeric → section number
FONT_TITLE = 20            # ≥ this → title
FONT_SUBTITLE = 11         # ≥ this → subtitle
FONT_ITEM_HEADER = 9.5     # bold around 10pt → item header
FONT_TAG_UPPER = 7.5       # ≤ this + uppercase + short → tag/label

# Recurring elements to filter out (matched by text content).
RECURRING_TEXTS = {
    "liderar con consciencia",
}

PAGE_NUM_RE = re.compile(r"^\d{1,3}\s*/\s*\d{1,3}$")

# Single-character bullets/decorators to strip (explicit set + general rule below)
BULLET_CHARS = {"•", "·", "●", "○", "►", "▪", "▸", "‣", "–", "—", "→", "➝", "➜",
                "!", "+", "↓", "↑", "←", "⬅", "↔", "⇒", "⇐", "⇔", "✓", "✗", "✔", "✕"}


def _is_decorative_symbol(text):
    """Return True if text is a single decorative symbol (arrows, bullets, PUA, etc.)."""
    if len(text) > 2:
        return False
    if text in BULLET_CHARS:
        return True
    # Filter single characters in Private Use Area (Font Awesome icons etc.)
    if len(text) == 1 and ord(text[0]) >= 0xE000:
        return True
    # Filter single non-alphanumeric, non-quote characters
    if len(text) == 1 and not text.isalnum() and text not in ('"', "'", '"', '"', '«', '»', '='):
        return True
    return False

# Source PPTX mapping: module key → filename
DIAPOS_DIR = Path(__file__).resolve().parent.parent / "curso" / "Presentación Powerpoint" / "Diapos Manus"

PPTX_MAP = {
    "M02": "M02-ppt-Consciente-como-estoy.pptx",
    "M03": "M03-ppt_Consciente_de_lo_que_pienso_y_siento.pptx",
    "M04": "M04-ppt-Consciente_de_lo_que_Necesito.pptx",
    "M05": "m05-ppt-Consciente_de_lo_que_necesitamos.pptx",
    "M06": "m06-ppt-Consciente_de_lo_que_Quiero.pptx",
    "M07": "m07-ppt-Consciente_de_lo_que_Está_Bien.pptx",
}

# Human-readable module names (for headers)
MODULE_NAMES = {
    "M02": "Consciente de Cómo Estoy",
    "M03": "Consciente de lo que Pienso y Siento",
    "M04": "Consciente de lo que Necesito",
    "M05": "Consciente de lo que Necesitamos",
    "M06": "Consciente de lo que Quiero",
    "M07": "Consciente de lo que Está Bien",
}

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def emu_to_inches(emu):
    return emu / EMU_PER_INCH if emu else 0

def emu_to_pt(emu):
    """Convert EMU to points (1 pt = 12700 EMU)."""
    return emu / 12700 if emu else 0


def get_dominant_font_size(shape):
    """Return the dominant font size in points for a shape's text.

    Walks all runs and picks the most common non-None size.
    Falls back to 0 if nothing found.
    """
    if not shape.has_text_frame:
        return 0
    sizes = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                sizes.append(run.font.size)
    if not sizes:
        return 0
    # Most common size
    from collections import Counter
    return emu_to_pt(Counter(sizes).most_common(1)[0][0])


def is_bold_dominant(shape):
    """Check if the dominant formatting is bold."""
    if not shape.has_text_frame:
        return False
    bold_chars = 0
    total_chars = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            n = len(run.text)
            total_chars += n
            if run.font.bold:
                bold_chars += n
    return total_chars > 0 and bold_chars / total_chars > 0.5


def shape_text(shape):
    """Get the full text of a shape, preserving paragraph breaks."""
    if not shape.has_text_frame:
        return ""
    lines = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs)
        lines.append(line)
    return "\n".join(lines).strip()


def shape_paragraphs(shape):
    """Get a list of paragraph texts with their run-level info."""
    if not shape.has_text_frame:
        return []
    result = []
    for para in shape.text_frame.paragraphs:
        text = "".join(run.text for run in para.runs)
        if not text.strip():
            continue
        # Get font info from first non-empty run
        font_size = 0
        bold = False
        for run in para.runs:
            if run.text.strip():
                if run.font.size:
                    font_size = emu_to_pt(run.font.size)
                bold = bool(run.font.bold)
                break
        result.append({
            "text": text.strip(),
            "font_size": font_size,
            "bold": bold,
        })
    return result


# ---------------------------------------------------------------------------
# Image classification
# ---------------------------------------------------------------------------

def classify_image(shape):
    """Classify a PICTURE shape by its dimensions and position."""
    w, h = shape.width, shape.height
    x, y = shape.left, shape.top

    # Full-slide background
    if abs(w - SLIDE_W_EMU) < BG_TOLERANCE and abs(h - SLIDE_H_EMU) < BG_TOLERANCE:
        return "full-bg"

    # Near-full background
    if w > 7000000 and h > 4000000:
        return "full-bg"

    # Small icon (< ~1" on both axes)
    if w < EMU_PER_INCH and h < EMU_PER_INCH:
        return "icon"

    # Decorative circle (approximately square, top-right area)
    if abs(w - h) < 200000 and x > SLIDE_W_EMU * 0.6 and y < SLIDE_H_EMU * 0.3:
        return "decorative"

    return "content"


# ---------------------------------------------------------------------------
# Shape extraction and filtering
# ---------------------------------------------------------------------------

def extract_shapes(slide):
    """Extract and classify all shapes from a slide.

    Returns (text_shapes, image_shapes) where each is a list of dicts.
    """
    text_shapes = []
    image_shapes = []

    for shape in slide.shapes:
        info = {
            "left": shape.left or 0,
            "top": shape.top or 0,
            "width": shape.width or 0,
            "height": shape.height or 0,
            "name": shape.name,
        }

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            info["classification"] = classify_image(shape)
            info["shape"] = shape
            image_shapes.append(info)
            continue

        if not shape.has_text_frame:
            continue

        text = shape_text(shape)
        if not text:
            continue

        info["text"] = text
        info["font_size"] = get_dominant_font_size(shape)
        info["bold"] = is_bold_dominant(shape)
        info["paragraphs"] = shape_paragraphs(shape)
        info["shape"] = shape
        text_shapes.append(info)

    return text_shapes, image_shapes


def filter_recurring(text_shapes):
    """Remove recurring/noise elements from text shapes."""
    filtered = []
    for s in text_shapes:
        text = s["text"]
        text_clean = text.strip()
        text_lower = text_clean.lower()

        # Header "Liderar con Consciencia"
        if text_lower in RECURRING_TEXTS:
            continue

        # Page number pattern "N / M"
        if PAGE_NUM_RE.match(text_clean):
            continue

        # Single decorative symbols (bullets, arrows, PUA icons)
        if _is_decorative_symbol(text_clean):
            continue

        # Small decorative slide numbers at the bottom of the slide
        # (1-2 digit numbers in the bottom zone with small font)
        if (s["top"] > BOTTOM_ZONE_Y
                and re.match(r"^\d{1,2}$", text_clean)
                and s["font_size"] < FONT_SUBTITLE):
            continue

        filtered.append(s)
    return filtered


# ---------------------------------------------------------------------------
# Role classification
# ---------------------------------------------------------------------------

ROLE_PRIORITY = [
    "section_number", "title", "subtitle", "tag",
    "stat", "quote", "item_header", "body",
]

def classify_role(s):
    """Assign a semantic role to a text shape based on font size and content."""
    text = s["text"]
    font = s["font_size"]
    bold = s["bold"]
    text_clean = text.strip()

    # Image placeholder
    if "[IMAGEN" in text or "[imagen" in text:
        return "image_placeholder"

    # Section number: large font, short numeric text
    if font >= FONT_SECTION_NUMBER and len(text_clean) <= 4 and re.match(r"^\d{1,2}$", text_clean):
        return "section_number"

    # Title: large font
    if font >= FONT_TITLE:
        return "title"

    # Subtitle: medium font
    if font >= FONT_SUBTITLE:
        return "subtitle"

    # Quote: text starts/ends with quotes or guillemets
    if (text_clean.startswith('"') or text_clean.startswith('"') or
            text_clean.startswith('«') or text_clean.startswith("'")):
        return "quote"

    # Tag/label: small font + uppercase + short
    if font <= FONT_TAG_UPPER and font > 0 and text_clean == text_clean.upper() and len(text_clean) < 50:
        return "tag"

    # Tag/label: small font + bold + short + in the top zone of the slide
    # (catches mixed-case labels like "Origen de las Emociones" at 7pt bold)
    if (font <= FONT_TAG_UPPER and font > 0 and bold
            and len(text_clean) < 40 and s["top"] < SLIDE_H_EMU * 0.25):
        return "tag"

    # Stat: short text with number/percentage as dominant element
    if (re.match(r"^[\d,.]+\s*[%xXseg]*$", text_clean) and len(text_clean) < 15
            and font >= 14):
        return "stat"

    # Item header: bold at ~10pt
    if bold and FONT_ITEM_HEADER <= font < FONT_SUBTITLE:
        return "item_header"

    # Source/attribution: very small font with academic/source feel
    if font > 0 and font <= 7 and len(text_clean) > 10:
        return "source"

    return "body"


# ---------------------------------------------------------------------------
# Reading-order grouping
# ---------------------------------------------------------------------------

# Vertical gap (EMU) between shapes to start a new visual group.
# Shapes separated by more than this gap vertically are in different groups.
GROUP_GAP_Y = EMU_PER_INCH * 0.35  # 0.35"

# Row band tolerance — shapes within this Y-distance are on the "same row"
ROW_BAND = EMU_PER_INCH * 0.20  # 0.20"


def _reading_order(shapes):
    """Sort shapes in reading order: top-to-bottom in row bands, left-to-right
    within each band.  Shapes whose top positions are within ROW_BAND of each
    other are considered the same visual row and sorted left-to-right."""
    if not shapes:
        return []

    # Sort primarily by top
    by_top = sorted(shapes, key=lambda s: s["top"])

    # Group into row bands
    rows = []
    current_row = [by_top[0]]
    for s in by_top[1:]:
        if s["top"] - current_row[0]["top"] <= ROW_BAND:
            current_row.append(s)
        else:
            rows.append(current_row)
            current_row = [s]
    rows.append(current_row)

    # Sort each row left-to-right and flatten
    ordered = []
    for row in rows:
        row.sort(key=lambda s: s["left"])
        ordered.extend(row)
    return ordered


def _group_by_vertical_gap(shapes):
    """Split an ordered list of shapes into groups separated by vertical gaps.
    Returns a list of lists."""
    if not shapes:
        return []

    groups = [[shapes[0]]]
    for s in shapes[1:]:
        prev = groups[-1][-1]
        prev_bottom = prev["top"] + prev["height"]
        gap = s["top"] - prev_bottom
        if gap > GROUP_GAP_Y:
            groups.append([s])
        else:
            groups[-1].append(s)
    return groups


# ---------------------------------------------------------------------------
# Markdown generation
# ---------------------------------------------------------------------------

def _clean_quotes(text):
    """Normalize curly/smart quotes to straight quotes, remove doubled quotes."""
    text = text.replace("\u201c", '"').replace("\u201d", '"')
    text = text.replace("\u2018", "'").replace("\u2019", "'")
    if text.startswith('""') and text.endswith('""'):
        text = '"' + text[2:-2] + '"'
    return text


def _format_shape(s):
    """Format a single shape as markdown based on its role."""
    role = s["role"]
    text = _clean_quotes(s["text"].strip())

    if role == "section_number":
        return f"**{text}**"
    if role == "tag":
        return f"`{text}`"
    if role == "title":
        return f"### {text}"
    if role == "subtitle":
        return f"*{text}*" if len(text) < 60 else text
    if role == "quote":
        if not (text.startswith('"') or text.startswith('«')):
            text = f'"{text}"'
        return f"> {text}"
    if role == "source":
        return f"— *{text}*"
    if role == "stat":
        return f"**{text}**"
    if role == "item_header":
        return f"- **{text}**"
    if role == "image_placeholder":
        return text  # already contains [IMAGEN: ...] markup
    # body
    return text if text else ""


def generate_slide_md(slide_num, text_shapes, image_shapes):
    """Generate markdown for one slide.

    Strategy: all shapes (text + images) are placed in **reading order**
    (row-banded top→bottom, left→right) and grouped by vertical proximity.
    Each group becomes a visual block separated by a blank line.
    """
    lines = []
    lines.append(f"## Diapositiva {slide_num}")
    lines.append("")

    # Classify roles
    for s in text_shapes:
        s["role"] = classify_role(s)
        s["kind"] = "text"

    # Content images (non-background, non-decorative)
    content_images = []
    for img in image_shapes:
        if img["classification"] == "content":
            w_in = emu_to_inches(img["width"])
            h_in = emu_to_inches(img["height"])
            img["role"] = "image"
            img["kind"] = "image"
            img["text"] = f"[IMAGEN {w_in:.1f}\"×{h_in:.1f}\"]"
            content_images.append(img)

    # Merge all items into a single list for reading-order sort
    all_items = list(text_shapes) + content_images
    if not all_items:
        lines.append("*(diapositiva vacía)*")
        lines.append("")
        lines.append("---")
        lines.append("")
        return "\n".join(lines)

    # Sort in reading order
    ordered = _reading_order(all_items)

    # Group by vertical gaps
    groups = _group_by_vertical_gap(ordered)

    # Render each group
    for gi, group in enumerate(groups):
        for s in group:
            if s["kind"] == "image":
                lines.append(s["text"])
            else:
                formatted = _format_shape(s)
                if formatted:
                    lines.append(formatted)
        # Blank line between groups (but not after the last one)
        if gi < len(groups) - 1:
            lines.append("")

    lines.append("")
    lines.append("---")
    lines.append("")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# PPTX processing
# ---------------------------------------------------------------------------

def process_pptx(module_key):
    """Process a single PPTX file and return the full markdown string."""
    filename = PPTX_MAP[module_key]
    filepath = DIAPOS_DIR / filename
    if not filepath.exists():
        print(f"ERROR: File not found: {filepath}", file=sys.stderr)
        return None

    prs = Presentation(str(filepath))
    slides = list(prs.slides)
    total = len(slides)
    name = MODULE_NAMES.get(module_key, module_key)

    md_lines = []
    md_lines.append(f"# {module_key} — {name}")
    md_lines.append("")
    md_lines.append(f"**Fuente:** {filename}")
    md_lines.append(f"**Diapositivas:** {total}")
    md_lines.append("")
    md_lines.append("---")
    md_lines.append("")

    for i, slide in enumerate(slides):
        slide_num = i + 1
        text_shapes, image_shapes = extract_shapes(slide)
        text_shapes = filter_recurring(text_shapes)

        slide_md = generate_slide_md(slide_num, text_shapes, image_shapes)
        md_lines.append(slide_md)

    return "\n".join(md_lines)


def write_output(module_key, md_content):
    """Write the markdown output file next to the source PPTX."""
    out_path = DIAPOS_DIR / f"{module_key}-slide-descriptions.md"
    out_path.write_text(md_content, encoding="utf-8")
    print(f"  -> {out_path.name}  ({len(md_content):,} chars)")
    return out_path


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Extract structured slide descriptions from Manus PPTX files."
    )
    parser.add_argument(
        "modules", nargs="*",
        help="Module keys to process (e.g. M02 M05). Use --all for all modules."
    )
    parser.add_argument(
        "--all", action="store_true",
        help="Process all 6 modules (M02-M07)."
    )
    args = parser.parse_args()

    if args.all:
        modules = list(PPTX_MAP.keys())
    elif args.modules:
        modules = [m.upper() for m in args.modules]
        for m in modules:
            if m not in PPTX_MAP:
                print(f"ERROR: Unknown module '{m}'. Valid: {', '.join(PPTX_MAP.keys())}", file=sys.stderr)
                sys.exit(1)
    else:
        parser.print_help()
        sys.exit(1)

    print(f"Processing {len(modules)} module(s): {', '.join(modules)}")
    print(f"Source dir: {DIAPOS_DIR}")
    print()

    for module_key in modules:
        print(f"[{module_key}] {PPTX_MAP[module_key]}")
        md = process_pptx(module_key)
        if md:
            write_output(module_key, md)

    print("\nDone.")


if __name__ == "__main__":
    main()
