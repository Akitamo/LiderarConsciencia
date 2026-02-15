"""
Compare two PPTX files slide by slide to find removed images/shapes.
Compares M04 original vs M04-PATRON.
"""
import sys
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

ORIGINAL = r"C:\dev\projects\LiderarConsciencia\_wip\curso\Presentación Powerpoint\ppts CLAUDE CODE\M04-Consciente-de-lo-que-Necesito.pptx"
PATRON   = r"C:\dev\projects\LiderarConsciencia\_wip\curso\Presentación Powerpoint\ppts CLAUDE CODE\M04-PATRON.pptx"

FULL_BG_W = 9144000   # standard 10" slide width in EMU
FULL_BG_H = 5143500   # standard 5.625" slide height in EMU
TOLERANCE = 50000     # ~0.05" tolerance for background detection

def classify_image(w, h):
    """Classify an image by its dimensions."""
    if abs(w - FULL_BG_W) < TOLERANCE and abs(h - FULL_BG_H) < TOLERANCE:
        return "FULL-SLIDE BACKGROUND"
    elif w > 7000000 and h > 4000000:
        return "NEAR-FULL BACKGROUND"
    elif (w > 8000000 and h < 500000) or (h > 4000000 and w < 500000):
        return "LINE/SEPARATOR (thin strip)"
    elif max(w, h) / max(min(w, h), 1) > 10:
        return "LINE/SEPARATOR (high aspect ratio)"
    elif w < 1000000 and h < 1000000:
        return "SMALL ICON/ELEMENT"
    else:
        return "DECORATIVE IMAGE"

def get_shape_info(slide):
    """Extract shape info from a slide."""
    pictures = []
    autoshapes = []
    textboxes = []
    other = []
    
    for shape in slide.shapes:
        info = {
            'name': shape.name,
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
            'shape_type': shape.shape_type,
        }
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # 13
            pictures.append(info)
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:  # 1
            autoshapes.append(info)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:  # 17
            textboxes.append(info)
        else:
            other.append(info)
    
    return pictures, autoshapes, textboxes, other

def emu_to_inches(emu):
    return emu / 914400

def main():
    prs_orig = Presentation(ORIGINAL)
    prs_patron = Presentation(PATRON)
    
    slides_orig = list(prs_orig.slides)
    slides_patron = list(prs_patron.slides)
    
    print(f"Original: {len(slides_orig)} slides")
    print(f"PATRON:   {len(slides_patron)} slides")
    print()
    
    max_slides = max(len(slides_orig), len(slides_patron))
    
    # Summary table header
    print("=" * 110)
    print(f"{'Slide':>5} | {'Pics Orig':>9} | {'Pics PAT':>8} | {'Diff':>5} | {'Auto Orig':>9} | {'Auto PAT':>8} | {'Diff':>5} | {'Text Orig':>9} | {'Text PAT':>8} | {'Diff':>5}")
    print("-" * 110)
    
    slides_with_removed_images = []
    
    for i in range(max_slides):
        slide_num = i + 1
        
        if i < len(slides_orig):
            pics_o, auto_o, text_o, other_o = get_shape_info(slides_orig[i])
        else:
            pics_o, auto_o, text_o, other_o = [], [], [], []
            
        if i < len(slides_patron):
            pics_p, auto_p, text_p, other_p = get_shape_info(slides_patron[i])
        else:
            pics_p, auto_p, text_p, other_p = [], [], [], []
        
        diff_pics = len(pics_o) - len(pics_p)
        diff_auto = len(auto_o) - len(auto_p)
        diff_text = len(text_o) - len(text_p)
        
        marker = ""
        if diff_pics > 0:
            marker = " <-- IMAGES REMOVED"
        
        print(f"{slide_num:>5} | {len(pics_o):>9} | {len(pics_p):>8} | {diff_pics:>+5} | {len(auto_o):>9} | {len(auto_p):>8} | {diff_auto:>+5} | {len(text_o):>9} | {len(text_p):>8} | {diff_text:>+5}{marker}")
        
        if diff_pics > 0:
            slides_with_removed_images.append({
                'num': slide_num,
                'orig_pics': pics_o,
                'patron_pics': pics_p,
            })
    
    print("=" * 110)
    
    # Detail for slides with removed images
    if slides_with_removed_images:
        print(f"\n\n{'='*80}")
        print("DETAIL: Slides where images were REMOVED")
        print(f"{'='*80}\n")
        
        lines_lost_slides = []
        
        for s in slides_with_removed_images:
            slide_num = s['num']
            orig = s['orig_pics']
            patron = s['patron_pics']
            
            # Find which images were removed by comparing dimensions
            patron_sigs = [(p['width'], p['height'], p['left'], p['top']) for p in patron]
            patron_matched = [False] * len(patron)
            
            removed = []
            kept = []
            
            for pic in orig:
                sig = (pic['width'], pic['height'], pic['left'], pic['top'])
                matched = False
                for j, ps in enumerate(patron_sigs):
                    if not patron_matched[j] and ps == sig:
                        patron_matched[j] = True
                        matched = True
                        kept.append(pic)
                        break
                if not matched:
                    removed.append(pic)
            
            print(f"--- Slide {slide_num} ---")
            print(f"  Original had {len(orig)} pictures, PATRON has {len(patron)} ({len(removed)} removed)\n")
            
            has_line_loss = False
            
            for pic in removed:
                w = pic['width']
                h = pic['height']
                classification = classify_image(w, h)
                w_in = emu_to_inches(w)
                h_in = emu_to_inches(h)
                
                is_full_bg = "FULL-SLIDE" in classification or "NEAR-FULL" in classification
                might_have_lines = not is_full_bg
                
                if might_have_lines:
                    has_line_loss = True
                
                print(f"  REMOVED: {pic['name']}")
                print(f"    Size: {w} x {h} EMU  ({w_in:.2f}\" x {h_in:.2f}\")")
                print(f"    Position: left={pic['left']}, top={pic['top']}")
                print(f"    Classification: {classification}")
                if is_full_bg:
                    print(f"    --> Likely just a solid/gradient background moved to master. NO lines lost.")
                else:
                    print(f"    --> *** POTENTIAL LINE/SEPARATOR LOSS -- needs visual check ***")
                print()
            
            if has_line_loss:
                lines_lost_slides.append(slide_num)
        
        # Final summary
        print(f"\n{'='*80}")
        print("SUMMARY: Potential line/separator losses")
        print(f"{'='*80}\n")
        
        if lines_lost_slides:
            print(f"Slides that may have lost lines/separators (non-full-background images removed):")
            for sn in lines_lost_slides:
                print(f"  - Slide {sn}")
            print(f"\nTotal: {len(lines_lost_slides)} slides need visual inspection")
        else:
            print("No slides appear to have lost lines or separators.")
            print("All removed images were full-slide backgrounds (likely moved to master slide).")
    else:
        print("\nNo images were removed in any slide.")

if __name__ == "__main__":
    main()
